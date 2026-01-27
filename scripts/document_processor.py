"""
Enhanced Document Processing Pipeline
Extracts text and maintains formatting from various document formats including
PDF, DOC/DOCX, TXT, MD, HTML, RTF, EPUB, and PPTX files.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any, Dict, List, Optional
import json

from pypdf import PdfReader  # type: ignore
import docx  # type: ignore
import markdown
from bs4 import BeautifulSoup
import pypandoc
from pptx import Presentation


class DocumentProcessor:
    """Extract and chunk textual documents with formatting preservation."""

    def __init__(self, chunk_size: int = 200, chunk_overlap: int = 40) -> None:
        if chunk_overlap >= chunk_size:
            raise ValueError("chunk_overlap must be smaller than chunk_size")

        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.doc': self._process_docx,
            '.docx': self._process_docx,
            '.txt': self._process_txt,
            '.md': self._process_markdown,
            '.html': self._process_html,
            '.htm': self._process_html,
            '.rtf': self._process_rtf,
            '.epub': self._process_epub,
            '.pptx': self._process_pptx,
        }

    def extract_metadata(self, text: str) -> Dict[str, Any]:
        """Extract metadata from document content."""
        metadata = {
            'word_count': len(text.split()),
            'char_count': len(text),
            'has_code': bool(re.search(r'```[\s\S]*?```|`.*?`', text)),
            'has_urls': bool(re.search(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', text)),
            'has_lists': bool(re.search(r'^\s*[-*+]\s|^\s*\d+\.\s', text, re.MULTILINE)),
            'has_tables': bool(re.search(r'\|.*\|.*\|', text) or '┌' in text or '╔' in text),
        }
        return metadata

    # ------------------------------------------------------------------
    def process_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Process a document file and return chunks with metadata."""
        ext = Path(file_path).suffix.lower()
        
        if ext not in self.supported_formats:
            raise ValueError(f"Unsupported document format: {ext}")
            
        processor = self.supported_formats[ext]
        return processor(file_path)
        
    def _process_markdown(self, file_path: str) -> List[Dict[str, Any]]:
        """Process Markdown files maintaining formatting."""
        filename = Path(file_path).name
        text = Path(file_path).read_text(encoding='utf-8', errors='ignore')
        
        # Convert Markdown to HTML to preserve formatting
        html = markdown.markdown(text, extensions=['tables', 'fenced_code'])
        soup = BeautifulSoup(html, 'html.parser')
        
        # Extract formatted text
        formatted_text = soup.get_text('\n', strip=True)
        
        return self._chunk_text(
            formatted_text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="markdown",
            file_path=file_path,
            extra_metadata={'source_format': 'markdown', 'has_formatting': True}
        )

    def _process_html(self, file_path: str) -> List[Dict[str, Any]]:
        """Process HTML files preserving structure."""
        filename = Path(file_path).name
        with open(file_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'html.parser')
        
        # Remove script and style elements
        for element in soup(['script', 'style']):
            element.decompose()
        
        # Extract text with basic formatting
        formatted_text = soup.get_text('\n', strip=True)
        
        return self._chunk_text(
            formatted_text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="html",
            file_path=file_path,
            extra_metadata={'source_format': 'html', 'has_formatting': True}
        )

    def _process_rtf(self, file_path: str) -> List[Dict[str, Any]]:
        """Process RTF files."""
        filename = Path(file_path).name
        
        # Convert RTF to HTML using pandoc
        html = pypandoc.convert_file(file_path, 'html', format='rtf')
        soup = BeautifulSoup(html, 'html.parser')
        formatted_text = soup.get_text('\n', strip=True)
        
        return self._chunk_text(
            formatted_text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="rtf",
            file_path=file_path,
            extra_metadata={'source_format': 'rtf', 'has_formatting': True}
        )

    def _process_epub(self, file_path: str) -> List[Dict[str, Any]]:
        """Process EPUB files."""
        filename = Path(file_path).name
        
        # Convert EPUB to HTML using pandoc
        html = pypandoc.convert_file(file_path, 'html', format='epub')
        soup = BeautifulSoup(html, 'html.parser')
        formatted_text = soup.get_text('\n', strip=True)
        
        return self._chunk_text(
            formatted_text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="epub",
            file_path=file_path,
            extra_metadata={'source_format': 'epub', 'has_formatting': True}
        )

    def _process_pptx(self, file_path: str) -> List[Dict[str, Any]]:
        """Process PowerPoint presentations."""
        filename = Path(file_path).name
        prs = Presentation(file_path)
        
        slides_text = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                formatted_text = '\n'.join(slide_text)
                chunks = self._chunk_text(
                    formatted_text,
                    filename=filename,
                    base_id=f"{filename}_s{idx}",
                    page=idx,
                    doc_type="pptx",
                    file_path=file_path,
                    extra_metadata={
                        'source_format': 'pptx',
                        'slide_number': idx,
                        'total_slides': len(prs.slides)
                    }
                )
                slides_text.extend(chunks)
        
        return slides_text

    # ------------------------------------------------------------------
    def _process_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        chunks: List[Dict[str, Any]] = []
        reader = PdfReader(file_path)
        filename = Path(file_path).name

        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                continue

            chunks.extend(
                self._chunk_text(
                    text,
                    filename=filename,
                    base_id=f"{filename}_p{page_number}",
                    page=page_number,
                    doc_type="pdf",
                    file_path=file_path,
                )
            )

        return chunks

    # ------------------------------------------------------------------
    def _process_docx(self, file_path: str) -> List[Dict[str, Any]]:
        document = docx.Document(file_path)
        text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
        filename = Path(file_path).name

        return self._chunk_text(
            text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="docx",
            file_path=file_path,
        )

    # ------------------------------------------------------------------
    def _process_txt(self, file_path: str) -> List[Dict[str, Any]]:
        filename = Path(file_path).name
        text = Path(file_path).read_text(encoding="utf-8", errors="ignore")

        return self._chunk_text(
            text,
            filename=filename,
            base_id=f"{filename}",
            page=None,
            doc_type="text",
            file_path=file_path,
        )

    # ------------------------------------------------------------------
    def _chunk_text(
        self,
        text: str,
        *,
        filename: str,
        base_id: str,
        page: int | None,
        doc_type: str,
        file_path: str,
        extra_metadata: Optional[Dict[str, Any]] = None,
    ) -> List[Dict[str, Any]]:
        words = text.split()
        if not words:
            return []

        chunks: List[Dict[str, Any]] = []
        step = self.chunk_size - self.chunk_overlap
        for start in range(0, len(words), step):
            end = min(start + self.chunk_size, len(words))
            content = " ".join(words[start:end]).strip()
            if not content:
                continue

            chunk_index = len(chunks)
            # Extract content-specific metadata
            content_metadata = self.extract_metadata(content)
            
            # Combine all metadata
            metadata = {
                "source": filename,
                "type": doc_type,
                "page": page,
                "file_path": file_path,
                "content_metadata": content_metadata,
            }
            
            # Add extra metadata if provided
            if extra_metadata:
                metadata.update(extra_metadata)
            
            chunks.append(
                {
                    "text": content,
                    "metadata": metadata,
                    "chunk_id": f"{base_id}_c{chunk_index}",
                }
            )

        return chunks


__all__ = ["DocumentProcessor"]
